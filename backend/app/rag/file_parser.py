import os
import uuid


class FileParser:
    SUPPORTED_TYPES = {
        "txt": "text",
        "md": "text",
        "py": "code",
        "js": "code",
        "ts": "code",
        "jsx": "code",
        "tsx": "code",
        "java": "code",
        "cpp": "code",
        "c": "code",
        "go": "code",
        "rs": "code",
        "rb": "code",
        "php": "code",
        "html": "code",
        "css": "code",
        "json": "code",
        "yaml": "code",
        "yml": "code",
        "xml": "code",
        "sql": "code",
        "sh": "code",
        "csv": "text",
        "pdf": "pdf",
        "docx": "docx",
        "xlsx": "xlsx",
        "pptx": "pptx",
        "png": "image",
        "jpg": "image",
        "jpeg": "image",
        "gif": "image",
        "webp": "image",
    }

    @classmethod
    def get_file_type(cls, filename: str) -> str:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        return cls.SUPPORTED_TYPES.get(ext, "unknown")

    @classmethod
    def extract_text(cls, file_path: str, file_type: str) -> str:
        if file_type == "text" or file_type == "code":
            return cls._read_text(file_path)
        elif file_type == "pdf":
            return cls._read_pdf(file_path)
        elif file_type == "docx":
            return cls._read_docx(file_path)
        elif file_type == "xlsx":
            return cls._read_xlsx(file_path)
        elif file_type == "pptx":
            return cls._read_pptx(file_path)
        elif file_type == "csv":
            return cls._read_csv(file_path)
        elif file_type == "image":
            return cls._ocr_image(file_path)
        return ""

    @staticmethod
    def _read_text(file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    @staticmethod
    def _read_pdf(file_path: str) -> str:
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        except ImportError:
            return "[PDF parsing requires pypdf]"

    @staticmethod
    def _read_docx(file_path: str) -> str:
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        except ImportError:
            return "[DOCX parsing requires python-docx]"

    @staticmethod
    def _read_xlsx(file_path: str) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            text = ""
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    text += " | ".join([str(c) if c else "" for c in row]) + "\n"
            return text
        except ImportError:
            return "[XLSX parsing requires openpyxl]"

    @staticmethod
    def _read_pptx(file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            return "[PPTX parsing requires python-pptx]"

    @staticmethod
    def _read_csv(file_path: str) -> str:
        import csv
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            return "\n".join([" | ".join(row) for row in reader])

    @staticmethod
    def _ocr_image(file_path: str) -> str:
        try:
            import pytesseract
            from PIL import Image
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        except ImportError:
            return "[Image OCR requires pytesseract and Pillow]"
